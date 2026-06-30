"""
Build both technical documents for the CleanMesh Studio project — v1.1
  - CleanMesh_기술문서.docx   (detailed Word manual)
  - CleanMesh_매뉴얼.pptx    (slide deck for presentation)

Reflects all features as of 2026-06-22:
  - D:\\GoogleDrive\\Image_to_3D project layout (실행/ + 작업리소스/ split)
  - SAM2 click-to-segment in WSL
  - Color modes: vertex / region_split
  - Color noise reduction (Laplacian smoothing)
  - K-means + label majority + tiny region merge
  - Input downsampling to 1024 px (avoid OOM)
  - Operating practices (close Blender + OBS during generation)
  - Deterministic mode (opt-in, 8 GB GPU OOM risk)

Run:
    python build_docs.py
"""

from __future__ import annotations

from pathlib import Path
from datetime import datetime

from docx import Document
from docx.shared import Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor as PPTRGBColor


HERE = Path(__file__).resolve().parent
OUT_DOCX = HERE / "CleanMesh_기술문서.docx"
OUT_PPTX = HERE / "CleanMesh_매뉴얼.pptx"


# =============================================================================
#   Word helpers
# =============================================================================
def _shade(cell, hex_color: str) -> None:
    tcPr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd"); shd.set(qn("w:fill"), hex_color); tcPr.append(shd)


def _set_kr_font(run, name="맑은 고딕"):
    run.font.name = name
    rPr = run._element.get_or_add_rPr()
    rFonts = rPr.find(qn("w:rFonts"))
    if rFonts is None:
        rFonts = OxmlElement("w:rFonts"); rPr.append(rFonts)
    rFonts.set(qn("w:eastAsia"), name); rFonts.set(qn("w:ascii"), name); rFonts.set(qn("w:hAnsi"), name)


def _para(doc, text, *, bold=False, size=10, color=None, align=None):
    p = doc.add_paragraph()
    if align is not None: p.alignment = align
    r = p.add_run(text); r.bold = bold; r.font.size = Pt(size)
    if color: r.font.color.rgb = RGBColor.from_string(color)
    _set_kr_font(r)


def _heading(doc, text, level=1):
    h = doc.add_heading(level=level); r = h.add_run(text); _set_kr_font(r)
    if level == 0: r.font.size = Pt(26); r.font.color.rgb = RGBColor.from_string("1f3864")
    elif level == 1: r.font.size = Pt(18); r.font.color.rgb = RGBColor.from_string("1f3864")
    elif level == 2: r.font.size = Pt(14); r.font.color.rgb = RGBColor.from_string("2e74b5")
    else: r.font.size = Pt(12)


def _code(doc, code):
    p = doc.add_paragraph()
    p.paragraph_format.left_indent = Cm(0.4); p.paragraph_format.right_indent = Cm(0.4)
    p.paragraph_format.space_before = Pt(4); p.paragraph_format.space_after = Pt(4)
    r = p.add_run(code); r.font.name = "Consolas"; r.font.size = Pt(9)
    r.font.color.rgb = RGBColor.from_string("0b3d91")
    pPr = p._element.get_or_add_pPr()
    shd = OxmlElement("w:shd"); shd.set(qn("w:fill"), "f4f4f7"); pPr.append(shd)


def _bullet(doc, text, *, level=0):
    p = doc.add_paragraph(style="List Bullet")
    p.paragraph_format.left_indent = Cm(0.6 + 0.5 * level)
    r = p.add_run(text); r.font.size = Pt(10); _set_kr_font(r)


def _table(doc, headers, rows, col_widths=None):
    t = doc.add_table(rows=1 + len(rows), cols=len(headers))
    t.alignment = WD_TABLE_ALIGNMENT.LEFT; t.style = "Light Grid Accent 1"
    for i, h in enumerate(headers):
        c = t.rows[0].cells[i]; c.text = ""
        r = c.paragraphs[0].add_run(h)
        r.bold = True; r.font.size = Pt(10); r.font.color.rgb = RGBColor.from_string("ffffff"); _set_kr_font(r)
        _shade(c, "1f3864")
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            c = t.rows[ri].cells[ci]; c.text = ""
            r = c.paragraphs[0].add_run(str(val)); r.font.size = Pt(9); _set_kr_font(r)
    if col_widths:
        for col, w in enumerate(col_widths):
            for row in t.rows: row.cells[col].width = Cm(w)


# =============================================================================
#   Word document
# =============================================================================
def build_docx():
    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "맑은 고딕"; style.font.size = Pt(10)
    rPr = style.element.get_or_add_rPr()
    rFonts = rPr.find(qn("w:rFonts"))
    if rFonts is None:
        rFonts = OxmlElement("w:rFonts"); rPr.append(rFonts)
    rFonts.set(qn("w:eastAsia"), "맑은 고딕")

    # Cover
    _heading(doc, "CleanMesh Studio", 0)
    _para(doc, "이미지/텍스트 → 디지털 트윈 3D 자동화 파이프라인", bold=True, size=14, color="2e74b5")
    _para(doc, "기술 문서 · 운영 매뉴얼 (v1.2 · Digital Twin)", size=11, color="595959")
    doc.add_paragraph()
    _table(doc, ["항목", "값"], [
        ["버전", "1.2 (Digital Twin)"],
        ["작성일", datetime.now().strftime("%Y-%m-%d")],
        ["프로젝트 루트", "D:\\GoogleDrive\\Image_to_3D\\"],
        ["대상 OS", "Windows 10/11 + WSL2 Ubuntu 22.04"],
        ["가속 장치", "NVIDIA GPU (CUDA 12.1, RTX 3070 8GB 검증, 16GB+ 권장)"],
        ["주요 의존성", "Blender 5.1.2 / Python 3.11 / PyTorch 2.4 / FastAPI / SAM2"],
        ["통신 포트", "8100 (CleanMesh API) · 8000 (TripoSR) · 9876 (Blender MCP)"],
    ], col_widths=[3.5, 12.5])
    doc.add_page_break()

    # 1. 개요
    _heading(doc, "1. 개요", 1)
    _para(doc,
          "CleanMesh Studio는 실제 산업 현장의 사진을 디지털 트윈 플랫폼"
          "(NVIDIA Omniverse, Twinmotion, BIM 협업툴, Pixar USD 기반 뷰어 등)에 "
          "임포트 가능한 정확도-우선 3D 자산으로 변환하는 엔드투엔드 파이프라인입니다. "
          "AGV·팔레트·드럼·컨베이어·선반 같은 현장 설비를 한 장의 사진과 마우스 클릭 한 "
          "번으로 디지털 트윈에 즉시 배치 가능한 GLB / USD 자산으로 만듭니다.")

    _heading(doc, "1.1 디지털 트윈 사용 시 핵심 가치", 2)
    _bullet(doc, "현장 설비 사진 → 실측 단위(mm) 메타데이터 포함 3D 자산 (1~2분)")
    _bullet(doc, "정확성 우선 정책: 정점 보존, 잡티만 최소 제거 (디테일 살림)")
    _bullet(doc, "vertex color로 실제 표면 외관 그대로 (산화·녹·라벨·페인트 자국)")
    _bullet(doc, "Watertight·manifold 메시 → 충돌·물리 시뮬레이션 가능")
    _bullet(doc, "GLB / USD / FBX 멀티 익스포트 — Omniverse · Twinmotion · BIM 협업툴 호환")
    _bullet(doc, "메타데이터: 카테고리·실치수·제조사·시리얼·생성시각 자동 임베드")
    _bullet(doc, "Procedural 표준 자산(EUR1 팔레트, 200L 드럼 등) 1초 생성")

    _heading(doc, "1.2 핵심 기능 (전체)", 2)
    _bullet(doc, "이미지 1장 → TRELLIS 고품질 (50~70초, vertex color 포함) — 권장")
    _bullet(doc, "이미지 1장 → TripoSR 빠른 프리뷰 (3~6초)")
    _bullet(doc, "이미지 2~3장 → TRELLIS Multi-view (1.5~2분)")
    _bullet(doc, "텍스트 → Procedural 템플릿 (드럼·팔레트·박스·셀프·컨베이어)")
    _bullet(doc, "SAM2 클릭 세그멘트로 배경 자동 제거 (ghost geometry 0)")
    _bullet(doc, "색상 모드 2종: 정점 색상 (DT 권장, 실표면 보존) / 영역별 분해 (자산 모듈화 용)")
    _bullet(doc, "잡티 제거 강도 4단계 — DT는 약 또는 끄기 권장 (디테일 보존)")
    _bullet(doc, "Blender 5.1 헤드리스 cleanup + 6방향 렌더 검증")
    _bullet(doc, "Blender MCP 연동: 결과물 원클릭 임포트")
    _bullet(doc, "입력 이미지 자동 다운샘플 (1024 px max) — VRAM 보호")
    _bullet(doc, "USD 익스포트 (Omniverse 호환) + GLB extras에 DT 메타데이터")

    _heading(doc, "1.3 처리량 벤치마크 (RTX 3070 8GB · 깨끗한 GPU 상태)", 2)
    _table(doc, ["경로", "입력", "총 소요시간", "정점 수", "비고"], [
        ["Procedural", "텍스트", "1~3초", "1,000~5,000", "Blender CLI"],
        ["TripoSR", "이미지 1장", "3~6초", "~50,000", "TripoSR FastAPI"],
        ["TRELLIS Single", "이미지 1장", "50~70초", "~250,000", "WSL2 subprocess"],
        ["TRELLIS Single + SAM2", "이미지 1장", "60~75초", "~250,000", "권장 패턴"],
        ["TRELLIS Multi (2~3장)", "이미지 2~3장", "1.5~2분", "~270,000", "VRAM 빠듯"],
        ["TRELLIS Multi (4장+)", "이미지 4장+", "—", "—", "❌ 8GB GPU OOM"],
        ["+ Color split (K=3)", "+K-means", "+10~15초", "정점 유지", "K개 머티리얼 분할"],
        ["+ Vertex smoothing", "+Laplacian", "+5~10초", "정점 유지", "잡티 제거"],
    ])

    # 2. 아키텍처
    doc.add_page_break()
    _heading(doc, "2. 시스템 아키텍처", 1)

    _heading(doc, "2.1 컴포넌트 구성", 2)
    _code(doc, """\
┌─────────────────────────────────────────────────────────────────────┐
│  Web Browser  →  http://localhost:8100  →  SPA + 캔버스             │
└───────────────────────────┬─────────────────────────────────────────┘
                            │ REST + multipart + JSON
┌───────────────────────────▼─────────────────────────────────────────┐
│  CleanMesh API (FastAPI, uvicorn, port 8100)                        │
│  ┌─────────────────────────────────────────────────────────────┐    │
│  │ server/main.py  · server/static/  (SAM2 캔버스 + 색상 옵션) │    │
│  └─────────────────────────────────────────────────────────────┘    │
│      │           │           │           │           │              │
│      ▼           ▼           ▼           ▼           ▼              │
│ Procedural   TripoSR    TRELLIS      SAM2       Color Process       │
│ (Blender     (port      (WSL2)       (WSL2)     (Blender headless)  │
│  --backg.)    8000)                                                 │
│      │           │           │           │           │              │
└──────┼───────────┼───────────┼───────────┼───────────┼──────────────┘
       ▼           ▼           ▼           ▼           ▼
  raw GLB → cleanup → [color_smooth or color_split] → render → export
                                                                  │
                                                                  ▼
                                                ┌────────────────┐
                                                │  Blender 5.1   │
                                                │  + MCP addon   │
                                                │  (port 9876)   │
                                                └────────────────┘
""")

    _heading(doc, "2.2 파이프라인 단계 (6 stage)", 2)
    _table(doc, ["단계", "역할", "구현 파일"], [
        ["1. Routing", "입력에 맞는 생성기 선택", "cleanmesh/router.py"],
        ["2. Generation", "Raw mesh 생성", "cleanmesh/generators/*"],
        ["3. Cleanup", "remove_doubles · Decimate · UV", "cleanmesh/blender/cleanup.py"],
        ["3.5. Color Process", "vertex smoothing 또는 region split", "cleanmesh/blender/color_*.py"],
        ["4. Render Verify", "6방향 PNG + contact sheet", "cleanmesh/blender/render.py"],
        ["5. Export", "엔진별 GLB/FBX", "cleanmesh/blender/export.py"],
    ], col_widths=[3.5, 7.5, 6.5])

    _heading(doc, "2.3 디렉토리 구조 (2026-06-22 기준)", 2)
    _code(doc, r"""
D:\GoogleDrive\Image_to_3D\                  ← 프로젝트 루트
├── CleanMesh_기술문서.docx                   ← 본 문서
├── CleanMesh_매뉴얼.pptx
├── build_docs.py                            ← 두 문서 재생성
├── README.md
│
├── 실행\                                     ← 바로 실행 가능
│   ├── CleanMeshStudio.exe                       ← PyInstaller 단일 .exe
│   ├── CleanMeshStudio.log
│   ├── CleanMesh_Start.bat / Stop.bat
│   ├── *.lnk (단축아이콘 사본)
│   ├── build_launcher.bat / migrate_pack.bat
│   └── _onedrive_disable_admin.bat / _onedrive_rollback_steps.txt
│
└── 작업리소스\                                 ← 소스 코드
    ├── cleanmesh_launcher.py
    ├── requirements.txt
    ├── cleanmesh\
    │   ├── router.py / pipeline.py / segment.py    ← segment.py (NEW v1.1)
    │   ├── generators\
    │   │   ├── procedural.py / triposr.py
    │   │   └── trellis.py                          ← --deterministic, 다운샘플 (NEW v1.1)
    │   └── blender\
    │       ├── cleanup.py / render.py / export.py
    │       ├── color_smooth.py                     ← NEW v1.1
    │       ├── color_split.py                      ← NEW v1.1
    │       └── templates\
    ├── server\ (main.py + static/index.html)
    └── trellis_wsl\                                ← WSL2 스크립트
        ├── _run_trellis.py / _sam2_segment.py
        └── _install_*.sh

C:\WorkingJob\3d-model-tool\python-backend\   ← TripoSR (별도)
C:\t3d\                                       ← 런타임 출력 (raw/cleaned/renders/exports/logs)
""")

    # 3. 환경
    doc.add_page_break()
    _heading(doc, "3. 환경 요구사항", 1)
    _heading(doc, "3.1 하드웨어", 2)
    _table(doc, ["부품", "최소", "권장", "운영 비고"], [
        ["GPU", "NVIDIA 8 GB VRAM", "16 GB+", "8GB는 멀티뷰 어려움 (8.3 참고)"],
        ["RAM", "16 GB", "32 GB", "WSL2 동적 할당"],
        ["디스크", "60 GB", "150 GB (D: 권장)", "모델 weights 17 GB"],
        ["CPU", "4 core", "8 core 이상", "Blender 헤드리스"],
    ])
    _heading(doc, "3.2 소프트웨어", 2)
    _table(doc, ["소프트웨어", "버전", "용도"], [
        ["Windows", "10/11", "호스트"],
        ["NVIDIA Driver", "591.86+", "CUDA 12.1"],
        ["Python (Windows)", "3.11.x", "API + 빌드"],
        ["WSL2 + Ubuntu", "22.04 LTS", "TRELLIS + SAM2"],
        ["Blender", "5.1.2", "메시 작업"],
        ["Blender MCP addon", "1.5.6", "port 9876"],
        ["TripoSR backend", "main", "단일 이미지"],
    ])

    # 4. 설치
    doc.add_page_break()
    _heading(doc, "4. 설치 가이드", 1)
    _heading(doc, "4.1 Windows 측", 2)
    _bullet(doc, "Python 3.11 설치 + PATH")
    _bullet(doc, "NVIDIA 드라이버 + CUDA 12.1")
    _bullet(doc, "Blender 5.1.2 설치")
    _bullet(doc, "프로젝트 의존성:")
    _code(doc, "cd 작업리소스\npython -m pip install -r requirements.txt")
    _bullet(doc, "Blender MCP addon: ahujasid/blender-mcp 의 zip → Blender preferences > Add-ons")
    _bullet(doc, ".exe 빌드: 실행/build_launcher.bat 더블클릭")

    _heading(doc, "4.2 TripoSR", 2)
    _code(doc, r"""
git clone https://github.com/VAST-AI-Research/TripoSR  C:\WorkingJob\3d-model-tool\python-backend
cd C:\WorkingJob\3d-model-tool\python-backend
python -m venv venv && venv\Scripts\activate
pip install -r requirements.txt""")

    _heading(doc, "4.3 WSL2 (TRELLIS + SAM2)", 2)
    _code(doc, r"""
wsl --install Ubuntu-22.04
# 그 후 trellis_wsl/_install_*.sh 순서대로:
bash /mnt/d/trellis/_install_apt.sh
bash /mnt/d/trellis/_install_step1.sh
bash /mnt/d/trellis/_install_step2.sh
bash /mnt/d/trellis/_install_kaolin.sh
bash /mnt/d/trellis/_fix_transformers.sh
bash /mnt/d/trellis/_sanity.sh
bash /mnt/d/trellis/_install_sam2_safe.sh    ← 반드시 _safe 버전""")
    _para(doc, "⚠ _install_sam2.sh (구버전) 실행 금지 — torch 2.4 → 2.12 강제 업그레이드로 TRELLIS 깨짐",
          bold=True, color="b22222")

    _heading(doc, "4.4 첫 실행", 2)
    _bullet(doc, "바탕화면 CleanMesh Studio 단축아이콘 더블클릭")
    _bullet(doc, "런처가 TripoSR(8000) + CleanMesh API(8100) 자동 기동")
    _bullet(doc, "브라우저가 자동으로 http://localhost:8100 열림")

    # 5. 모듈별
    doc.add_page_break()
    _heading(doc, "5. 모듈별 상세", 1)

    _heading(doc, "5.1 Router", 2)
    _table(doc, ["조건", "선택", "비고"], [
        ["text only", "Procedural", "이미지 없음"],
        ["image × 1, quality=speed", "TripoSR", "빠른 프리뷰"],
        ["image × 1, quality=quality", "TRELLIS Single", "고품질"],
        ["image × 2~3", "TRELLIS Multi", "멀티뷰 융합"],
        ["image × 4+", "TRELLIS Multi", "⚠ 8GB GPU OOM 위험"],
    ])

    _heading(doc, "5.2 Procedural 템플릿", 2)
    _bullet(doc, "drum_200l, pallet_eur, box_cargo, shelf_rack, conveyor_roller")

    _heading(doc, "5.3 TripoSR", 2)
    _code(doc, """\
POST http://localhost:8000/generate
  file: <PNG/JPG>
  remove_bg: true | false   ← SAM2 사용 시 false""")

    _heading(doc, "5.4 TRELLIS", 2)
    _para(doc, "Windows측에서 입력 이미지를 1024 px max 다운샘플 + ASCII 이름으로 복사 후 "
               "WSL의 _run_trellis.py 를 subprocess 호출.")
    _code(doc, r"""
wsl -d Ubuntu-22.04 -- bash -lc "source ~/trellis-venv/bin/activate && \
  python /mnt/d/trellis/_run_trellis.py \
    --inputs <list> --output <out.glb> \
    --seed 1 --steps-ss 12 --steps-slat 12 \
    [--pre-masked] [--deterministic]"

Vertex color: 3D Gaussian DC term → KNN 으로 정점에 전이""")

    _heading(doc, "5.5 SAM2 클릭 세그멘트 (NEW v1.1)", 2)
    _para(doc, "Meta SAM 2.1 hiera_large (224M 파라미터). 1~5개 클릭으로 정밀 마스크 → "
               "RGBA로 TRELLIS에 전달 → TRELLIS rembg 우회.")
    _bullet(doc, "양성 클릭 (label=1): 주체 위, 녹색 점")
    _bullet(doc, "음성 클릭 (label=0): 배경 위, 빨강 점 (Shift+클릭)")
    _bullet(doc, "출력: RGBA PNG 컷아웃 + 마스크 + bbox + confidence")
    _bullet(doc, "feather=2px 가장자리 부드럽게")
    _table(doc, ["메서드", "경로", "설명"], [
        ["POST", "/api/segment-upload", "이미지 N장 → session_id"],
        ["POST", "/api/segment", "클릭 좌표 → 마스크 PNG"],
        ["GET", "/api/segment/{sid}/{file}", "세션 파일 서빙"],
        ["POST", "/api/generate-segmented", "세션의 마스크로 TRELLIS 시작"],
    ])

    _heading(doc, "5.6 색상 처리 2모드 (NEW v1.1)", 2)
    _table(doc, ["모드", "처리", "결과", "용도"], [
        ["vertex (DT 권장)",     "정점마다 색상 + Laplacian 평탄화", "실제 표면 그대로",   "디지털 트윈 (산화·녹·라벨 보존)"],
        ["region_split",        "K-means → K개 머티리얼",       "평평한 단색 K개",   "자산 모듈화 / 부품별 머티리얼 교체 / LOD"],
    ])
    _para(doc, "디지털 트윈에서는 vertex 모드를 강력 권장합니다. 실제 설비의 산화·녹·"
               "라벨·페인트 자국이 색상 그대로 보존되어 현장 식별성과 진단 가치가 유지됩니다.",
          bold=True, color="2e74b5")
    _para(doc, "두 모드 공통 잡티 제거 강도:")
    _table(doc, ["강도", "smooth_iters", "label_smooth_iters", "min_region_size"], [
        ["강", "7", "4", "200"],
        ["중 (기본)", "5", "3", "100"],
        ["약", "2", "1", "30"],
        ["끄기", "0", "0", "1"],
    ])

    _heading(doc, "5.6.1 Vertex smoothing", 3)
    _para(doc, "메시 edge 따라 Laplacian smoothing — 인접 정점들의 색상 평균. "
               "TRELLIS KNN 전이 후 점박이 노이즈가 부드럽게. numpy.add.at 벡터화로 295k 정점이 ~1초.")

    _heading(doc, "5.6.2 Region split", 3)
    _bullet(doc, "① 정점 색상 Laplacian smoothing")
    _bullet(doc, "② LAB 색공간 K-means (K-means++ 시드)")
    _bullet(doc, "③ 라벨 majority filter (1-ring 이웃 다수결)")
    _bullet(doc, "④ 면별 majority vote (3 정점 다수결)")
    _bullet(doc, "⑤ 작은 connected component 흡수 (Union-Find)")
    _bullet(doc, "⑥ K개 머티리얼 생성 (cluster center 색)")
    _bullet(doc, "결과: 단일 GLB 안에 K개 primitive → 엔진에서 K개 sub-mesh + K개 머티리얼")

    _heading(doc, "5.7 Blender Cleanup", 2)
    _bullet(doc, "중복 정점 머지 (threshold 0.0001)")
    _bullet(doc, "Manifold 검사 + 갭 채우기")
    _bullet(doc, "Decimate Collapse → target_polys 도달")
    _bullet(doc, "Smart UV unwrap (66°)")
    _bullet(doc, "법선 재계산 (auto-smooth 35°)")
    _bullet(doc, "Vertex color → Base Color 자동 연결")

    _heading(doc, "5.8 Render Verification", 2)
    _para(doc, "EEVEE Next 6방향 (정면·후면·좌·우·위·아래) 256×256 PNG + 3×2 contact_sheet.png")

    # 6. API
    doc.add_page_break()
    _heading(doc, "6. REST API", 1)
    _heading(doc, "6.1 헬스/메타", 2)
    _table(doc, ["메서드", "경로", "설명"], [
        ["GET", "/api/health", "전체 컴포넌트 상태"],
        ["GET", "/api/health/blender", "Blender MCP 연결 + 씬"],
        ["GET", "/api/templates", "Procedural 템플릿 목록"],
    ])
    _heading(doc, "6.2 생성", 2)
    _table(doc, ["메서드", "경로", "주요 인자"], [
        ["POST", "/api/generate", "file, quality, color_mode, color_k, color_smooth_iters"],
        ["POST", "/api/generate-multi", "files[], color_mode, color_k"],
        ["POST", "/api/generate-procedural", "template, params(JSON)"],
        ["POST", "/api/generate-segmented", "session_id, image_indices, color_mode, color_k"],
    ])
    _heading(doc, "6.3 SAM2", 2)
    _table(doc, ["메서드", "경로", "설명"], [
        ["POST", "/api/segment-upload", "이미지 → session"],
        ["POST", "/api/segment", "클릭 → 마스크"],
        ["GET", "/api/segment/{sid}/{file}", "세션 파일"],
    ])
    _heading(doc, "6.4 상태/결과", 2)
    _table(doc, ["메서드", "경로", "설명"], [
        ["GET", "/api/status/{job_id}", "잡 상태 폴링"],
        ["GET", "/api/download/{job_id}", "GLB 다운로드"],
        ["GET", "/api/contact-sheet/{job_id}", "6방향 합성"],
        ["POST", "/api/import-to-blender", "Blender로 push"],
    ])

    # 7. UI
    doc.add_page_break()
    _heading(doc, "7. 웹 UI 사용법", 1)
    _heading(doc, "7.1 일반 흐름", 2)
    _bullet(doc, "단축아이콘 더블클릭 → 브라우저 자동 오픈")
    _bullet(doc, "이미지 드래그&드롭 + 옵션 선택 + ▶ Generate")
    _bullet(doc, "결과 카드 + contact sheet + 다운로드/Blender 버튼")

    _heading(doc, "7.2 SAM2 클릭", 2)
    _bullet(doc, "🎯 체크 → 캔버스에 이미지 → 주체 좌클릭 → ✓ 마스크 생성")
    _bullet(doc, "여러 장은 탭으로 전환 후 각각 마스크")

    _heading(doc, "7.3 색상 옵션", 2)
    _bullet(doc, "라디오: 정점 색상 / 영역별 분해 (K=2~8 슬라이더)")
    _bullet(doc, "잡티 제거 강도: 강/중/약/끄기")

    _heading(doc, "7.4 Blender 연동", 2)
    _bullet(doc, "Blender 5.1 + MCP addon 가 port 9876 listen")
    _bullet(doc, "헤더 상태 녹색 확인 후 🎨 Blender로 보내기")
    _bullet(doc, "⚠ Generate 시작 전 Blender 닫아야 안전 (8.3 참고)")

    # 8. 디버깅
    doc.add_page_break()
    _heading(doc, "8. 디버깅 가이드", 1)

    _heading(doc, "8.1 로그 위치", 2)
    _table(doc, ["로그", "경로", "내용"], [
        ["런처", r"D:\GoogleDrive\Image_to_3D\실행\CleanMeshStudio.log", "포트 / PID / 배너"],
        ["API", r"C:\t3d\logs\cleanmesh.log", "uvicorn + pipeline"],
        ["TripoSR", r"C:\t3d\logs\triposr.log", "단일 이미지 변환"],
        ["TRELLIS", "WSL stdout → cleanmesh.log", "RESULT:{...}"],
        ["SAM2", "WSL stderr → cleanmesh.log", "model load + clicks"],
    ])

    _heading(doc, "8.2 자주 발생하는 문제", 2)
    _table(doc, ["증상", "원인", "조치"], [
        ["바닥/벽이 같이 솟음", "rembg 한계", "SAM2 클릭으로 마스크"],
        ["TRELLIS import error", "torch 2.4 깨짐", "_install_sam2_safe.sh 재실행"],
        ["Blender 상태 빨강", "MCP addon 미실행", "MCP > Start Server"],
        ["session not found", "서버 재시작", "이미지 재업로드"],
        ["같은 사진인데 결과 다름", "디퓨전 비결정성", "--deterministic (단 OOM 위험)"],
        ["CUDA out of memory", "다른 GPU 앱 점유", "8.3 운영 수칙 참고"],
    ])

    _heading(doc, "8.3 ★ GPU 메모리 운영 수칙 (8 GB VRAM) ★", 2)
    _para(doc, "RTX 3070 8 GB 환경에서 다른 GPU 점유 앱이 켜져있으면 거의 100% OOM. 실측:",
          bold=True, color="b22222")
    _table(doc, ["조합", "결과"], [
        ["깨끗한 환경 + 단일이미지", "✅ 정상 (60~70초)"],
        ["깨끗한 환경 + 2~3장 멀티뷰", "⚠ 빠듯하지만 성공"],
        ["깨끗한 환경 + 4장+ 멀티뷰", "❌ OOM"],
        ["Blender 열린 채로 + 단일", "❌ OOM"],
        ["Blender 열린 채로 + 멀티", "❌ OOM + WSL 강제 종료"],
        ["OBS/Teams/Discord 켠 채로", "❌ VRAM 단편화"],
    ])
    _para(doc, "원인:")
    _bullet(doc, "WDDM은 25+개 프로세스에 GPU 컨텍스트 share → VRAM 단편화")
    _bullet(doc, "TRELLIS는 mesh extraction에서 1~2 GB 연속 블록 필요")
    _bullet(doc, "다른 앱이 작은 chunk를 사방에 잡으면 연속 블록 못 얻음")
    _bullet(doc, "WSL VM(vmmem) 죽으면 작업폴더 휘발")
    _para(doc, "권장 순서:", bold=True)
    _bullet(doc, "① Blender / OBS / Teams / Discord / NVIDIA Share / Razer / Chrome 다탭 닫기")
    _bullet(doc, "② CleanMesh Studio 실행 → 이미지 업로드 → SAM2 → Generate")
    _bullet(doc, "③ ✅ 완료 확인")
    _bullet(doc, "④ Blender 열기 → MCP addon 자동 시작 대기")
    _bullet(doc, "⑤ 🎨 Blender로 보내기 클릭 → import")

    _heading(doc, "8.4 진단 명령어", 2)
    _code(doc, r"""
Get-NetTCPConnection -LocalPort 8100,8000,9876 -State Listen     # 포트
nvidia-smi                                                       # GPU + 프로세스
wsl -d Ubuntu-22.04 -- bash -lc "free -h; nvidia-smi"            # WSL 상태
Get-Content "D:\GoogleDrive\Image_to_3D\실행\CleanMeshStudio.log" -Wait -Tail 30

# 잡 시작 전 GPU 정리
Stop-Process -Name obs64,Teams,Discord,RazerAppEngine,steamwebhelper -Force -ErrorAction SilentlyContinue""")

    # 9. 권장 패턴
    doc.add_page_break()
    _heading(doc, "9. 디지털 트윈 사용 권장 패턴", 1)

    _heading(doc, "9.1 디지털 트윈 자산 제작 표준 워크플로우", 2)
    _bullet(doc, "① 현장 사진 1장 촬영 (정면 또는 자산 식별성 가장 높은 각도)")
    _bullet(doc, "② 사진을 웹 UI에 드래그&드롭")
    _bullet(doc, "③ 🎯 SAM2 체크 → 주체 클릭 → ✓ 마스크 생성")
    _bullet(doc, "④ 색상 모드: 정점 색상 (vertex) — 표면 외관 보존")
    _bullet(doc, "⑤ 잡티 제거 강도: 약 또는 끄기 — 디테일 우선")
    _bullet(doc, "⑥ 메타데이터 입력: 카테고리·실치수(mm)·제조사·시리얼")
    _bullet(doc, "⑦ 익스포트 포맷: USD (Omniverse) 또는 GLB (범용)")
    _bullet(doc, "⑧ ▶ Generate → ~70초")
    _bullet(doc, "⑨ Blender 열기 → MCP로 import → 디지털 트윈에 배치")

    _heading(doc, "9.2 디지털 트윈 권장 옵션 표", 2)
    _para(doc, "정확성·진단가치·재현성을 우선하는 디지털 트윈 기본 설정:")
    _table(doc, ["항목", "권장 값", "이유"], [
        ["품질",          "quality (TRELLIS Single)",   "DINOv2 feature + 24 step diffusion 으로 디테일 보존"],
        ["색상 모드",      "vertex (정점 색상)",          "실제 표면 외관(산화·녹·라벨) 그대로 보존 → 자산 식별 가치"],
        ["잡티 제거 강도",  "약 (light, 기본값)",         "Laplacian 2회만 — 디퓨전 노이즈는 제거하되 디테일은 유지"],
        ["target_polys", "None (보존)",                "감산(decimate) 없이 원본 정점 그대로 → 정확성 우선"],
        ["익스포트 포맷",   "USD (Omniverse) 또는 GLB",   "USD: customLayerData 임베드 / GLB: extras 임베드"],
        ["메타데이터",     "카테고리·실치수·제조사·시리얼", "DT 플랫폼에서 자산 검색·필터링·진단에 사용"],
        ["엔진 옵션",      "omniverse / twinmotion / bim", "실세계 단위(meter) 보존, Z-up USD 컨벤션"],
        ["Deterministic", "16GB+ GPU 시 ON",            "같은 사진 → 같은 모델 → 자산 버전 관리·diff 가능"],
    ])

    _heading(doc, "9.3 8 GB GPU 사용 시 제한사항", 2)
    _table(doc, ["우선순위", "패턴", "추천도"], [
        ["1", "단일 이미지 + SAM2 + vertex color", "⭐⭐⭐ DT 표준"],
        ["2", "Procedural 템플릿 (드럼/팔레트/박스 등)", "⭐⭐⭐ 항상 동작"],
        ["3", "2~3장 멀티뷰 + SAM2", "⭐⭐ 360° 자산 필요시"],
        ["4", "4장+ 멀티뷰", "❌ 8GB GPU OOM"],
    ])

    _heading(doc, "9.4 16+ GB GPU 환경", 2)
    _bullet(doc, "5~7장 멀티뷰로 완전 360° 자산 안정적 생성")
    _bullet(doc, "Blender 열어둔 채로 generate 가능 (DT 작업 흐름 끊기지 않음)")
    _bullet(doc, "--deterministic 항상 활성화 → 같은 사진 → 같은 자산 (재현 가능 DT)")

    _heading(doc, "9.5 Deterministic mode (재현성)", 2)
    _para(doc, "디지털 트윈에서는 같은 설비를 다시 스캔했을 때 같은 모델이 나와야 "
               "버전 관리·diff 비교가 가능합니다. --deterministic 옵션이 이를 보장합니다.")
    _bullet(doc, "장점: 같은 입력 + 같은 SAM2 클릭 → byte-near-identical (정점 변동 ±0.01%)")
    _bullet(doc, "단점: 30% 느려짐 + 8 GB GPU에선 mesh extraction OOM 위험")
    _bullet(doc, "default OFF (8 GB 안정성 우선) — 16 GB+ GPU에선 ON 권장")

    _heading(doc, "9.6 메타데이터 임베드 규칙", 2)
    _para(doc, "GLB extras 및 USD attributes에 자동 임베드되는 필드:")
    _table(doc, ["필드", "예시 값", "용도"], [
        ["category",           "AGV / pallet / drum / shelf",   "자산 분류"],
        ["dimensions_mm",      "[1200, 800, 144]",              "실측 치수 (W×D×H)"],
        ["manufacturer",       "현대모비스 / Toyota",            "제조사"],
        ["serial_number",      "AGV-2024-001",                  "자산 고유 ID"],
        ["scan_timestamp",     "2026-06-23T14:30:00",           "스캔 시각"],
        ["source_image",       "IMG_0295.png",                  "원본 사진 파일명"],
        ["pipeline_version",   "CleanMesh v1.2",                "생성 시스템 버전"],
    ])

    # ─── 9.7 SimReady ──────────────────────────────────────────
    _heading(doc, "9.7 SimReady USD 출력 (NVIDIA Omniverse 호환)", 2)
    _para(doc,
          "USD 익스포트시 'SimReady' 옵션을 켜면 OpenUSD 기반 NVIDIA SimReady "
          "스펙에 호환되는 자산이 됩니다. 일반 USD는 정적 지오메트리만 담지만, "
          "SimReady 자산은 Omniverse / Twinmotion / IsaacSim 에서 즉시 "
          "물리 시뮬레이션이 가능합니다. 본 시스템은 NVIDIA Omniverse Kit 설치 "
          "없이 pxr (USD Python) 만으로 SimReady 호환 USD를 생성합니다.")
    _table(doc, ["주입 항목", "USD 스키마", "값 결정 방식"], [
        ["Stage upAxis",         "UsdGeom.SetStageUpAxis",       "Y 강제"],
        ["Stage metersPerUnit",  "UsdGeom.SetStageMetersPerUnit","1.0 (meters)"],
        ["defaultPrim",          "stage.SetDefaultPrim",         "최상위 Xform 자동 지정"],
        ["Kind",                 "Usd.ModelAPI.SetKind",         "component"],
        ["assetInfo",            "Usd.ModelAPI.SetAssetInfo",    "name, version, manufacturer, serial_number"],
        ["Real-world scale",     "UsdGeom.Xformable.AddScaleOp", "dimensions_mm 입력시 자동 보정"],
        ["PhysicsCollisionAPI",  "UsdPhysics.CollisionAPI",      "모든 메시에 적용"],
        ["MeshCollisionAPI",     "approximation=convexHull",     "8GB GPU 적합 (TriangleMesh 대비 가벼움)"],
        ["MassAPI",              "UsdPhysics.MassAPI",           "카테고리별 density (steel 7800·wood 700·aluminum 2700·box 200 kg/m³)"],
        ["PhysicsMaterial",      "UsdPhysics.MaterialAPI",       "μs=0.5 / μd=0.5 / e=0.0 (default)"],
        ["semanticLabel",        "Prim.SetCustomDataByKey",      "dt_meta.category에서 자동"],
    ])
    _para(doc,
          "검증: --simready 옵션으로 출력 후 시스템이 자동으로 SimReady "
          "셀프체크를 수행합니다. NVIDIA aif-pipeline (Omniverse Kit 필요) 이 "
          "PATH에 있으면 공식 validator로 검증하고, 없으면 내장 pxr 기반 "
          "체크리스트로 대체합니다. 두 경우 모두 검증 실패가 파이프라인 "
          "성공 여부에 영향을 주지 않습니다 — 진단 정보 용도입니다.")

    # 10. 마이그레이션
    doc.add_page_break()
    _heading(doc, "10. 다른 PC로 이전", 1)
    _para(doc, "실행/migrate_pack.bat → 3가지 모드:")
    _table(doc, ["모드", "번들", "셋업 시간", "추천 상황"], [
        ["A: 코드만", "~10 MB", "3~4시간", "인터넷 빠름"],
        ["B: 코드 + WSL", "~30 GB", "30분", "★ 일반 추천"],
        ["C: 풀", "~50 GB", "15분", "동일 GPU 모델"],
    ])

    # 11. 부록
    doc.add_page_break()
    _heading(doc, "부록 A. 치트시트", 1)
    _code(doc, r"""
# 서비스
바탕화면 > CleanMesh Studio                       # 시작
.\실행\CleanMesh_Stop.bat                         # 종료

# 빌드
.\실행\build_launcher.bat
.\실행\migrate_pack.bat [D:\대상]

# 디버깅
nvidia-smi
curl http://localhost:8100/api/health
type C:\t3d\logs\cleanmesh.log

# WSL
wsl --list --verbose
wsl --export Ubuntu-22.04 D:\bk.tar

# GPU 정리 (잡 시작 전)
Stop-Process -Name obs64,Teams,Discord,RazerAppEngine,steamwebhelper -Force""")

    _heading(doc, "부록 B. 라이센스", 1)
    _bullet(doc, "TripoSR: MIT (VAST-AI-Research)")
    _bullet(doc, "TRELLIS: MIT (Microsoft)")
    _bullet(doc, "SAM 2.1: Apache 2.0 (Meta)")
    _bullet(doc, "Blender: GPL v2/v3")
    _bullet(doc, "blender-mcp: MIT (ahujasid)")

    _heading(doc, "문서 정보", 1)
    _para(doc, f"버전 1.1 · {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    _para(doc, "재생성: python build_docs.py")

    doc.save(OUT_DOCX)
    print(f"OK  Word: {OUT_DOCX}")


# =============================================================================
#   PowerPoint
# =============================================================================
C_BG    = PPTRGBColor(0x1a, 0x1c, 0x20)
C_PANEL = PPTRGBColor(0x24, 0x27, 0x2d)
C_TEXT  = PPTRGBColor(0xe8, 0xea, 0xed)
C_MUTED = PPTRGBColor(0x9a, 0xa0, 0xa6)
C_ACC   = PPTRGBColor(0x4f, 0x8c, 0xff)
C_OK    = PPTRGBColor(0x34, 0xc7, 0x59)
C_WARN  = PPTRGBColor(0xff, 0x95, 0x00)
C_WHITE = PPTRGBColor(0xff, 0xff, 0xff)


def _slide_bg(slide, color):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = color


def _tb(slide, l, t, w, h, *, text, font="맑은 고딕", size=18, bold=False,
        color=C_TEXT, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = PPt(size); r.font.bold = bold; r.font.color.rgb = color
    return tb


def _line(slide, l, t, w, color, weight_pt=2):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Emu(int(weight_pt * 12700)))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()


def _bullets(slide, l, t, w, h, items, *, size=14, color=C_TEXT, indent_color=C_ACC):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(it, tuple): label, body = it
        else: label, body = "•", it
        r1 = p.add_run(); r1.text = label + "  "
        r1.font.name = "맑은 고딕"; r1.font.size = PPt(size); r1.font.bold = True; r1.font.color.rgb = indent_color
        r2 = p.add_run(); r2.text = body
        r2.font.name = "맑은 고딕"; r2.font.size = PPt(size); r2.font.color.rgb = color


def _p_table(slide, l, t, w, h, headers, rows):
    ts = slide.shapes.add_table(1 + len(rows), len(headers), l, t, w, h)
    table = ts.table
    for i, header in enumerate(headers):
        c = table.cell(0, i)
        c.fill.solid(); c.fill.fore_color.rgb = C_ACC
        p = c.text_frame.paragraphs[0]
        r = p.add_run(); r.text = header
        r.font.name = "맑은 고딕"; r.font.size = PPt(13); r.font.bold = True; r.font.color.rgb = C_WHITE
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            c = table.cell(ri, ci)
            c.fill.solid(); c.fill.fore_color.rgb = C_PANEL
            p = c.text_frame.paragraphs[0]
            r = p.add_run(); r.text = str(val)
            r.font.name = "맑은 고딕"; r.font.size = PPt(11); r.font.color.rgb = C_TEXT


def _p_code(slide, l, t, w, h, code, *, size=10):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.fill.solid(); tb.fill.fore_color.rgb = PPTRGBColor(0x0e, 0x10, 0x14)
    tb.line.fill.background()
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line
        r.font.name = "Consolas"; r.font.size = PPt(size)
        r.font.color.rgb = PPTRGBColor(0xb4, 0xc8, 0xff)


def _title(slide, title, subtitle=None):
    _tb(slide, PInches(0.5), PInches(0.3), PInches(12), PInches(0.7),
        text=title, size=28, bold=True, color=C_WHITE)
    _line(slide, PInches(0.5), PInches(1.0), PInches(12), C_ACC, weight_pt=2)
    if subtitle:
        _tb(slide, PInches(0.5), PInches(1.05), PInches(12), PInches(0.4),
            text=subtitle, size=12, color=C_MUTED)


def _footer(slide, idx, total):
    _tb(slide, PInches(0.5), PInches(7.0), PInches(12), PInches(0.3),
        text=f"CleanMesh Studio v1.2 (Digital Twin)  ·  {idx}/{total}",
        size=9, color=C_MUTED, align=PP_ALIGN.RIGHT)


def build_pptx():
    prs = Presentation()
    prs.slide_width = PInches(13.333); prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]
    slides = []

    def new():
        s = prs.slides.add_slide(blank); _slide_bg(s, C_BG); slides.append(s); return s

    # 1. Title
    s = new()
    _tb(s, PInches(0.5), PInches(2.5), PInches(12), PInches(1.2),
        text="🛠 CleanMesh Studio", size=54, bold=True, color=C_WHITE)
    _tb(s, PInches(0.5), PInches(3.7), PInches(12), PInches(0.6),
        text="이미지·텍스트 → 디지털 트윈 3D 자동화 파이프라인",
        size=20, color=C_ACC)
    _line(s, PInches(0.5), PInches(4.5), PInches(4), C_ACC, weight_pt=3)
    _tb(s, PInches(0.5), PInches(4.7), PInches(12), PInches(0.5),
        text="v1.2 · 디지털 트윈 자산 제작 자동화", size=14, color=C_MUTED)
    _tb(s, PInches(0.5), PInches(6.5), PInches(12), PInches(0.4),
        text=datetime.now().strftime("%Y-%m-%d"), size=11, color=C_MUTED)

    # 2. 시스템 한눈에
    s = new(); _title(s, "1. 시스템 한눈에", "디지털 트윈 자산을 사진 한 장으로")
    _bullets(s, PInches(0.6), PInches(1.6), PInches(12), PInches(5), [
        ("①", "현장 설비(AGV·팔레트·드럼 등) 사진 1장 촬영"),
        ("②", "마우스로 주체를 클릭 (SAM2가 정밀 마스크)"),
        ("③", "AI (TRELLIS) 가 정점 색상 포함 3D 메시 생성"),
        ("④", "Blender 자동 정리 + 표면 외관 그대로 보존"),
        ("⑤", "메타데이터(치수·제조사·시리얼) 자동 임베드"),
        ("⑥", "USD / GLB로 익스포트 → Omniverse / Twinmotion / BIM 협업툴로"),
    ], size=16)

    # 3. 생성 경로
    s = new(); _title(s, "2. 4가지 생성 경로", "입력에 따라 자동 선택")
    _p_table(s, PInches(0.5), PInches(1.6), PInches(12.3), PInches(4),
           ["경로", "입력", "소요시간", "특징"],
           [["Procedural", "텍스트", "1~3초", "정밀 산업 모델"],
            ["TripoSR", "이미지 1장", "3~6초", "빠른 프리뷰"],
            ["TRELLIS Single", "이미지 1장", "~60초", "고품질, vertex color"],
            ["TRELLIS Multi", "이미지 2~3장", "~2분", "부분 360° (8GB GPU 한계)"]])
    _tb(s, PInches(0.5), PInches(6.0), PInches(12), PInches(0.5),
        text="※ SAM2 클릭 세그멘트를 켜면 모든 경로에서 배경 ghost geometry 제거",
        size=12, color=C_OK)

    # 4. 아키텍처
    s = new(); _title(s, "3. 아키텍처")
    _p_code(s, PInches(0.5), PInches(1.6), PInches(12.3), PInches(5.2), """\
   Browser  (HTTP)
       │
       ▼
  ┌─────────────────────────────────────────────────┐
  │   CleanMesh API  (FastAPI, port 8100)           │
  │   • /api/generate*           작업 큐             │
  │   • /api/segment*            SAM2 클릭           │
  │   • /api/import-to-blender   MCP 푸시            │
  └─────────────────────────────────────────────────┘
       │           │           │           │
       ▼           ▼           ▼           ▼
  Procedural   TripoSR    TRELLIS (WSL)   SAM2 (WSL)
                                          hiera_large
       │           │           │           │
       ▼           ▼           ▼           ▼
              raw GLB → cleanup → color → render → export
                                                     │
                                                     ▼
                                       Blender 5.1 (MCP port 9876)""")

    # 5. 하드웨어
    s = new(); _title(s, "4. 하드웨어 요구사항")
    _p_table(s, PInches(0.5), PInches(1.6), PInches(12.3), PInches(4.5),
           ["부품", "최소", "권장", "비고"],
           [["GPU", "NVIDIA 8 GB", "16 GB+", "8GB는 단일 위주"],
            ["RAM", "16 GB", "32 GB", "WSL2 동적 할당"],
            ["디스크", "60 GB", "150 GB", "모델 weights 17 GB"],
            ["CPU", "4 core", "8 core+", "Blender 헤드리스"]])
    _tb(s, PInches(0.5), PInches(6.2), PInches(12), PInches(0.5),
        text="검증 환경: RTX 3070 8 GB + 16 GB RAM + WSL2 Ubuntu 22.04",
        size=12, color=C_MUTED)

    # 6. 1분 사용법
    s = new(); _title(s, "5. 1분 사용법", "더블클릭 한 번이면 끝")
    _bullets(s, PInches(0.6), PInches(1.6), PInches(12), PInches(5), [
        ("1.", "바탕화면 \"CleanMesh Studio\" 더블클릭"),
        ("2.", "브라우저 자동 오픈 → http://localhost:8100"),
        ("3.", "이미지 끌어다 놓기 (1장 권장)"),
        ("4.", "🎯 SAM2 체크 → 주체 클릭 → ✓ 마스크 생성"),
        ("5.", "색상 옵션 선택 (정점 / 영역 분해)"),
        ("6.", "▶ Generate → 60~70초 대기"),
        ("7.", "Blender 켜기 → 🎨 Blender로 보내기"),
    ], size=16)

    # 7. SAM2
    s = new(); _title(s, "6. SAM2 클릭 세그멘트", "배경 ghost geometry 완전 제거")
    _bullets(s, PInches(0.6), PInches(1.6), PInches(12), PInches(5), [
        ("✓", "옵션의 \"🎯 클릭으로 주체만 추출 (SAM2)\" 체크"),
        ("✓", "주체 위에 좌클릭 (녹색 점)"),
        ("✓", "배경 잡히면 Shift+클릭 (빨강 점)"),
        ("✓", "✓ 마스크 생성 → 1~2초 후 컷아웃 미리보기"),
        ("✓", "여러 이미지면 탭으로 다른 이미지 선택 후 반복"),
        ("✓", "▶ Generate → ghost geometry 0의 깨끗한 메시"),
    ], size=16)
    _tb(s, PInches(0.5), PInches(6.4), PInches(12), PInches(0.5),
        text="→ 단일 이미지 + SAM2 한 번이면 공장 바닥/벽 100% 제거",
        size=12, color=C_OK)

    # 8. 색상 2모드
    s = new(); _title(s, "7. 색상 처리 — DT는 vertex 모드 권장")
    _p_table(s, PInches(0.5), PInches(1.6), PInches(12.3), PInches(2.5),
           ["모드", "결과", "디지털 트윈 적합도"],
           [["정점 색상 (vertex)", "실제 표면 그라데이션 그대로", "⭐⭐⭐ 권장 — 산화·녹·라벨 보존"],
            ["영역별 분해 (region_split)", "K개 sub-mesh, 평평한 단색", "단색 LOD · 모듈화 용 (DT 비권장)"]])
    _bullets(s, PInches(0.6), PInches(4.4), PInches(12), PInches(2.5), [
        ("→", "DT에서는 자산 진단·식별을 위해 실제 외관 보존이 중요"),
        ("→", "잡티 제거 강도: 약 또는 끄기 권장 (디테일 우선)"),
        ("→", "vertex color는 GLB·USD 모두에서 자동 PBR 변환"),
    ], size=14)

    # 9. 잡티 제거 강도
    s = new(); _title(s, "8. 잡티 제거 강도 (v1.1)", "Laplacian smoothing + 라벨 정리")
    _p_table(s, PInches(0.5), PInches(1.6), PInches(12.3), PInches(3.5),
           ["강도", "정점 smooth", "라벨 smooth", "min region", "효과"],
           [["강", "7회", "4회", "200면", "잡티 최대 제거"],
            ["중 (기본)", "5회", "3회", "100면", "권장 — 잡티/디테일 균형"],
            ["약", "2회", "1회", "30면", "디테일 보존"],
            ["끄기", "—", "—", "—", "원본 그대로"]])
    _tb(s, PInches(0.5), PInches(5.6), PInches(12), PInches(0.5),
        text="vertex 모드 = smoothing만 / region_split = smoothing + K-means + 패치 흡수",
        size=12, color=C_MUTED)

    # 10. 벤치마크
    s = new(); _title(s, "9. 처리량 벤치마크", "RTX 3070 8 GB · 깨끗한 GPU 상태")
    _p_table(s, PInches(0.5), PInches(1.6), PInches(12.3), PInches(4.5),
           ["경로", "입력", "총 시간", "출력 크기"],
           [["Procedural", "텍스트", "1~3초", "0.3 MB"],
            ["TripoSR", "이미지 1장", "3~6초", "~5 MB"],
            ["TRELLIS Single + SAM2", "이미지 1장", "60~70초", "~18 MB"],
            ["TRELLIS Multi (2~3장)", "이미지 2~3장", "~2분", "~19 MB"],
            ["+ Color split (K=3)", "+K-means", "+15초", "~15 MB"],
            ["+ Vertex smoothing", "+Laplacian", "+8초", "변동 없음"]])

    # 11. 출력 위치
    s = new(); _title(s, "10. 출력 파일 위치")
    _p_table(s, PInches(0.5), PInches(1.6), PInches(12.3), PInches(4.5),
           ["종류", "경로", "포맷"],
           [["원본", "C:\\t3d\\raw\\", "GLB"],
            ["정리됨", "C:\\t3d\\cleaned\\", "GLB (_split / _smooth 변종)"],
            ["6방향 렌더", "C:\\t3d\\renders\\…\\", "PNG + contact_sheet.png"],
            ["엔진별", "C:\\t3d\\exports\\", "GLB"],
            ["SAM2 세션", "C:\\t3d\\raw\\segment_sessions\\{id}\\", "PNG"],
            ["로그", "C:\\t3d\\logs\\", "*.log"]])

    # 12. GPU 메모리 운영 수칙 (NEW)
    s = new(); _title(s, "11. ⚠ GPU 메모리 운영 수칙", "8 GB GPU에서 OOM 안 나려면")
    _tb(s, PInches(0.5), PInches(1.5), PInches(12), PInches(0.5),
        text="잡 시작 전에 반드시 닫아야 하는 앱들:", size=15, bold=True, color=C_WARN)
    _bullets(s, PInches(0.6), PInches(2.1), PInches(12), PInches(4.5), [
        ("❌", "Blender (잡 끝나고 켜기)"),
        ("❌", "OBS Studio"),
        ("❌", "Microsoft Teams"),
        ("❌", "Discord"),
        ("❌", "NVIDIA GeForce Experience / Share"),
        ("❌", "Razer Synapse"),
        ("❌", "Chrome 다른 탭들 (CleanMesh만 남김)"),
        ("✓", "권장: 생성 → ✅ 완료 → Blender 켜기 → Blender로 보내기"),
    ], size=14)

    # 13. 자주 발생 문제
    s = new(); _title(s, "12. 자주 발생하는 문제 Top 5")
    _p_table(s, PInches(0.5), PInches(1.6), PInches(12.3), PInches(5),
           ["증상", "원인", "해결"],
           [["CUDA out of memory", "다른 GPU 앱 점유", "Blender/OBS/Teams 닫기"],
            ["바닥/벽이 같이 솟음", "rembg 한계", "SAM2 클릭"],
            ["같은 사진인데 결과 다름", "디퓨전 비결정성", "--deterministic (8GB OOM 위험)"],
            ["TRELLIS import error", "torch 2.4 깨짐", "_install_sam2_safe.sh 재실행"],
            ["Blender 상태 빨강", "MCP 미실행", "MCP > Start Server"]])

    # 14. API
    s = new(); _title(s, "13. REST API 핵심 엔드포인트", "http://localhost:8100")
    _p_code(s, PInches(0.5), PInches(1.6), PInches(12.3), PInches(5.2), """\
GET  /api/health                  → 컴포넌트 상태
GET  /api/health/blender          → MCP 연결

POST /api/generate                → 이미지 1장 → 3D
POST /api/generate-multi          → 이미지 N장
POST /api/generate-procedural     → 템플릿
POST /api/generate-segmented      → SAM2 컷아웃으로 3D
    body: color_mode, color_k, color_smooth_iters, ...

POST /api/segment-upload          → 이미지 → session_id
POST /api/segment                 → 클릭 → 마스크 PNG
GET  /api/segment/{sid}/{file}    → 세션 파일

GET  /api/status/{job_id}         → 폴링
GET  /api/download/{job_id}       → GLB
GET  /api/contact-sheet/{job_id}  → 6방향 합성
POST /api/import-to-blender       → Blender push""")

    # 15. 디렉토리
    s = new(); _title(s, "14. 디렉토리 구조 (v1.1)", "D:\\GoogleDrive\\Image_to_3D")
    _p_code(s, PInches(0.5), PInches(1.6), PInches(12.3), PInches(5.2), r"""
D:\GoogleDrive\Image_to_3D\
├── CleanMesh_기술문서.docx
├── CleanMesh_매뉴얼.pptx
├── build_docs.py
├── README.md
│
├── 실행\
│   ├── CleanMeshStudio.exe       ← PyInstaller .exe
│   ├── CleanMesh_Start.bat / Stop.bat
│   ├── build_launcher.bat / migrate_pack.bat
│   └── *.lnk (단축아이콘 사본)
│
└── 작업리소스\
    ├── cleanmesh_launcher.py
    ├── requirements.txt
    ├── cleanmesh\
    │   ├── router.py / pipeline.py / segment.py
    │   ├── generators\  (procedural / triposr / trellis)
    │   └── blender\
    │       (cleanup / render / export / color_smooth / color_split)
    ├── server\  (main.py + static/index.html)
    └── trellis_wsl\
        (_run_trellis.py / _sam2_segment.py / _install_*.sh)""")

    # 16. 마이그레이션
    s = new(); _title(s, "15. 다른 PC로 이전")
    _p_table(s, PInches(0.5), PInches(1.6), PInches(12.3), PInches(3.5),
           ["모드", "번들 크기", "셋업 시간", "추천 상황"],
           [["A: 코드만", "~10 MB", "3~4시간", "인터넷 빠름"],
            ["B: 코드 + WSL", "~30 GB", "30분", "★ 추천"],
            ["C: 풀", "~50 GB", "15분", "동일 GPU"]])
    _bullets(s, PInches(0.6), PInches(5.3), PInches(12), PInches(1.5), [
        ("✓", "번들과 함께 setup_on_new_pc.bat 자동 생성"),
        ("✓", "robocopy + wsl --import + .exe 재빌드까지 자동"),
    ], size=14)

    # 17. Deterministic
    s = new(); _title(s, "16. Deterministic mode", "재현성 vs VRAM")
    _p_table(s, PInches(0.5), PInches(1.6), PInches(12.3), PInches(2.5),
           ["옵션", "정점 변동", "시간", "안정성 (8GB GPU)"],
           [["default OFF", "±15%", "60초", "✅ 안정"],
            ["--deterministic", "±0.01%", "70초", "❌ mesh extraction OOM"]])
    _bullets(s, PInches(0.6), PInches(4.4), PInches(12), PInches(2.5), [
        ("→", "TRELLIS는 디퓨전 → 같은 seed로도 매번 다른 결과 (본질)"),
        ("→", "cuDNN deterministic + 시드 고정으로 byte-near-identical 가능"),
        ("→", "하지만 8 GB GPU 에선 추가 메모리 비용으로 OOM"),
        ("→", "16+ GB GPU 사용자만 활성화 권장"),
    ], size=14)

    # 18. 향후 확장
    s = new(); _title(s, "17. 향후 확장 후보")
    _bullets(s, PInches(0.6), PInches(1.6), PInches(12), PInches(5), [
        ("□", "Grounded-SAM: \"빨간 박스만\" 텍스트 → 자동 마스크"),
        ("□", "멀티뷰 SAM2 일괄: 한 클릭이 모든 뷰에 자동 매핑"),
        ("□", "Camera projection 텍스처링: 라벨/디테일 보존"),
        ("□", "잡 영속화 (SQLite/Redis)"),
        ("□", "LOD 자동 생성 (target_polys 여러 단계 동시)"),
        ("□", "Web UI 모바일 대응"),
        ("□", "VRAM 자동 보호: 잡 시작 전 다른 앱 자동 종료/경고"),
    ], size=14)

    # 18.5 디지털 트윈 워크플로우 (NEW)
    s = new(); _title(s, "18. 디지털 트윈 표준 워크플로우", "현장 사진 → DT 자산 1분")
    _bullets(s, PInches(0.6), PInches(1.6), PInches(12), PInches(5), [
        ("1", "현장 설비 사진 1장 촬영 (자산 식별성 가장 높은 각도)"),
        ("2", "🎯 SAM2 체크 → 주체 클릭 → ✓ 마스크 생성"),
        ("3", "색상: vertex / 잡티 제거: 약 / 익스포트: USD"),
        ("4", "메타데이터 입력 — 카테고리, 실치수(mm), 제조사, 시리얼"),
        ("5", "☑ SimReady 자산으로 출력 체크 (USDPhysics 자동 주입)"),
        ("6", "▶ Generate → 약 70초"),
        ("7", "Omniverse / Twinmotion / BIM 협업툴에서 import → 시뮬레이션 즉시 가능"),
        ("→", "재현성: --deterministic 켜면 같은 사진 → 같은 모델 (16GB+ GPU)"),
    ], size=14)

    # 18.7 SimReady 슬라이드 (NEW)
    s = new(); _title(s, "19. SimReady USD 출력", "NVIDIA Omniverse 호환 디지털 트윈 자산")
    _p_table(s, PInches(0.5), PInches(1.6), PInches(12.3), PInches(3.6),
           ["항목", "일반 USD", "SimReady USD (이 시스템)"],
           [["메시 + 정점 색상",          "✅", "✅"],
            ["upAxis=Y / metersPerUnit=1.0", "종종 누락", "✅ 강제"],
            ["Kind=component + assetInfo",  "✗", "✅"],
            ["USDPhysics Collider (convex hull)",  "✗", "✅ 모든 메시"],
            ["USDPhysics Mass (카테고리별 density)", "✗", "✅ steel/wood/Al/box 자동"],
            ["USDPhysics Material (마찰·반발)", "✗", "✅ 기본값 주입"],
            ["semanticLabel + 실측 mm 스케일", "✗", "✅"]])
    _bullets(s, PInches(0.6), PInches(5.5), PInches(12), PInches(2), [
        ("→", "Omniverse Kit 설치 불필요 — pxr (Blender 5.x 번들 USD Python) 만 사용"),
        ("→", "옵션: aif-pipeline 설치된 환경에선 NVIDIA 공식 validator 자동 호출"),
    ], size=13)

    # 19. 요약
    s = new()
    _tb(s, PInches(0.5), PInches(2.5), PInches(12), PInches(1.0),
        text="✓ 한 줄 요약", size=32, bold=True, color=C_OK)
    _tb(s, PInches(0.5), PInches(3.7), PInches(12), PInches(1.0),
        text="현장 설비 사진 1장으로, 60초 안에 디지털 트윈에 배치 가능한 3D 자산이 나옵니다.",
        size=22, color=C_WHITE)
    _tb(s, PInches(0.5), PInches(4.8), PInches(12), PInches(0.6),
        text="vertex color + 약한 잡티 제거 + USD 익스포트 = DT 표준 출력",
        size=16, color=C_MUTED)
    _line(s, PInches(0.5), PInches(5.8), PInches(4), C_ACC, weight_pt=3)
    _tb(s, PInches(0.5), PInches(6.0), PInches(12), PInches(0.5),
        text="자세한 내용 → CleanMesh_기술문서.docx (v1.2)",
        size=14, color=C_MUTED)

    total = len(slides)
    for i, sl in enumerate(slides, 1):
        if i == 1: continue
        _footer(sl, i, total)

    prs.save(OUT_PPTX)
    print(f"OK  PPT:  {OUT_PPTX}  ({total} slides)")


if __name__ == "__main__":
    build_docx()
    build_pptx()
    print("\nDone.")
